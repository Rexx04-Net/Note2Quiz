import os
import sys
import json
import uuid
import datetime
import time
import re
import random   
import string   
import google.generativeai as genai
from flask import Flask, request, jsonify
from flask_cors import CORS
from pymongo import MongoClient
from pypdf import PdfReader
from pptx import Presentation
from youtube_transcript_api import YouTubeTranscriptApi

# --- CONFIGURATION ---
print("🔍 --- STARTING BACKEND ---")
current_dir = os.path.dirname(os.path.abspath(__file__))
GEMINI_API_KEY = None

# Secure Key Loading
try:
    sys.path.append(current_dir)
    from api_secrets import GEMINI_API_KEY
    print("✅ Successfully imported API Key.")
except:
    pass

app = Flask(__name__)
CORS(app, resources={r"/*": {"origins": "*"}})

# --- DATABASE CONNECTION ---
USING_MONGO = False
try:
    client = MongoClient("mongodb://localhost:27017/", serverSelectionTimeoutMS=2000)
    db = client["note2quiz_db"]
    notebooks_col = db["notebooks"]
    users_col = db["users"]  
    logs_col = db["activity_logs"] # ✅ NEW: Dedicated folder for tracking actions!
    
    client.server_info()
    print("✅ Connected to MongoDB")
    USING_MONGO = True
except:
    print("⚠️ MongoDB not found. Falling back to memory storage.")
    memory_notebooks = []
    
# Active Games Memory
active_games = {}

# --- AI CONFIGURATION ---
if GEMINI_API_KEY:
    genai.configure(api_key=GEMINI_API_KEY)

MODEL_PRIORITY = [
    'gemini-2.5-flash',
    'models/gemini-2.5-flash',
    'gemini-1.5-flash',
    'gemini-2.0-flash',
]

def generate_with_fallback(prompt):
    last_error = None
    for model_name in MODEL_PRIORITY:
        try:
            print(f"🤖 Trying AI Model: {model_name}...")
            model = genai.GenerativeModel(model_name)
            response = model.generate_content(prompt)
            print(f"✅ Success with {model_name}")
            return response.text
        except Exception as e:
            print(f"⚠️ Failed with {model_name}: {e}")
            last_error = e
            if "400" in str(e) or "API_KEY" in str(e): break
            time.sleep(1)
    
    return f"AI Error: All models failed. Last error: {str(last_error)}"

# ✅ NEW: HELPER FUNCTION FOR TRACKING ACTIVITY
def log_activity(email, action, details=""):
    if USING_MONGO:
        try:
            log_entry = {
                "email": email,
                "action": action,
                "details": details,
                "timestamp": str(datetime.datetime.now())
            }
            logs_col.insert_one(log_entry)
            print(f"📝 LOGGED: {email} -> {action}")
        except Exception as e:
            print(f"⚠️ Failed to log activity: {e}")

# --- HELPER FUNCTIONS ---
def get_notebook(notebook_id):
    if USING_MONGO:
        return notebooks_col.find_one({"id": notebook_id})
    else:
        return next((n for n in memory_notebooks if n["id"] == notebook_id), None)

def save_notebook(notebook):
    if USING_MONGO:
        notebooks_col.replace_one({"id": notebook["id"]}, notebook, upsert=True)
    else:
        for i, n in enumerate(memory_notebooks):
            if n["id"] == notebook["id"]:
                memory_notebooks[i] = notebook
                return
        memory_notebooks.append(notebook)

def get_context(notebook_id):
    nb = get_notebook(notebook_id)
    if not nb: return ""
    context = ""
    for source in nb.get("sources", []):
        context += f"\n--- SOURCE: {source['title']} ({source['type']}) ---\n{source['content']}\n"
    return context

def clean_ai_response(text):
    try: return json.loads(text)
    except:
        match = re.search(r'(\[.*\]|\{.*\})', text, re.DOTALL)
        if match:
            try: return json.loads(match.group(0))
            except: pass
    return None

def extract_text_from_pdf(file_stream):
    try:
        reader = PdfReader(file_stream)
        return "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
    except: return ""

def extract_text_from_pptx(file_stream):
    try:
        prs = Presentation(file_stream)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"): text.append(shape.text)
        return "\n".join(text)
    except: return ""

# --- ENDPOINTS ---
@app.route('/', methods=['GET'])
def home(): return "Note2Quiz Backend is Running! 🚀"

@app.route('/login', methods=['POST'])
def login():
    data = request.json
    email = data.get('email', '').strip().lower()
    
    if not email:
        return jsonify({"error": "Email is required"}), 400

    if USING_MONGO:
        existing_user = users_col.find_one({"email": email})
        
        if not existing_user:
            new_user = {
                "email": email,
                "created_at": str(datetime.datetime.now()),
                "role": "student" 
            }
            users_col.insert_one(new_user)
            log_activity(email, "Registered Account", "First time user sign-in") # ✅ Logging Registration
        else:
            log_activity(email, "Logged In", "User signed in successfully") # ✅ Logging normal Login

    return jsonify({"success": True, "email": email, "message": "Login successful"})

@app.route('/create-notebook', methods=['POST'])
def create_notebook():
    data = request.json
    email = data.get('email', 'guest')
    title = data.get('title', 'Untitled Notebook')
    
    new_nb = {
        "id": str(uuid.uuid4()),
        "user_email": email,
        "title": title,
        "sources": [],
        "created_at": str(datetime.datetime.now())
    }
    save_notebook(new_nb)
    
    # ✅ Logging Notebook Creation
    log_activity(email, "Created Notebook", f"Created notebook titled: {title}")
    
    if USING_MONGO: new_nb.pop('_id', None)
    return jsonify(new_nb)

@app.route('/get-notebooks', methods=['POST'])
def get_notebooks():
    email = request.json.get('email', 'guest')
    if USING_MONGO:
        notebooks = list(notebooks_col.find({"user_email": email}, {"_id": 0}))
    else:
        notebooks = [n for n in memory_notebooks if n['user_email'] == email]
    return jsonify(notebooks)

@app.route('/add-source', methods=['POST'])
def add_source():
    notebook_id = request.form.get('notebook_id')
    source_type = request.form.get('type', 'text')
    content = ""
    title = "New Source"

    if 'file' in request.files:
        f = request.files['file']
        title = f.filename
        if f.filename.endswith('.pdf'): content = extract_text_from_pdf(f)
        elif f.filename.endswith('.pptx'): content = extract_text_from_pptx(f)
    else:
        content = request.form.get('content', '')

    if not content.strip(): return jsonify({"error": "No text extracted"}), 400

    new_source = {
        "id": str(uuid.uuid4()),
        "title": title,
        "content": content,
        "type": source_type,
        "date": str(datetime.datetime.now())
    }
    
    nb = get_notebook(notebook_id)
    if nb:
        nb.setdefault('sources', []).append(new_source)
        save_notebook(nb)
        
        # ✅ Logging Source Upload
        user_email = nb.get('user_email', 'unknown')
        log_activity(user_email, "Added Source", f"Added {source_type} source: {title}")
        
        return jsonify(new_source)
    return jsonify({"error": "Notebook not found"}), 404

@app.route('/generate-studio-item', methods=['POST'])
def generate_studio_item():
    data = request.json
    notebook_id = data.get('notebook_id')
    tool_type = data.get('tool_type')
    difficulty = data.get('difficulty', 'Standard')
    num_questions = 5 if difficulty == "Easy" else 20 if difficulty == "Hard" else 10

    # ✅ Logging AI Generation Start
    nb = get_notebook(notebook_id)
    user_email = nb.get('user_email', 'unknown') if nb else 'unknown'
    log_activity(user_email, f"Generated {tool_type.capitalize()}", f"Triggered Gemini AI for {difficulty} {tool_type}")

    context = get_context(notebook_id)
    if len(context) < 50: return jsonify({"type": "text", "data": "Not enough content. Add sources first."})

    prompt = ""
    if tool_type == "quiz":
        prompt = f"""Generate a {difficulty} Quiz with exactly {num_questions} questions.
        STRICT JSON FORMAT:
        [
            {{
                "question": "Question text here?",
                "options": ["Option A", "Option B", "Option C", "Option D"],
                "answer": "Option A",
                "hint": "A short, helpful clue without giving the answer away."
            }}
        ]
        Context: {context[:35000]}"""

    elif tool_type == "flashcard":
        prompt = f"""Generate 10 Flashcards. Format: [{{ "front": "Term", "back": "Definition" }}]. Context: {context[:35000]}"""

    elif tool_type == "mindmap":
        prompt = f"""Create a hierarchical Mind Map using Emojis (🌳, 🌿). No Markdown blocks. Context: {context[:35000]}"""

    elif tool_type == "report":
        prompt = f"""Write an Executive Briefing Doc with bold headers. Context: {context[:35000]}"""

    try:
        text = generate_with_fallback(prompt)
        if text.startswith("AI Error:"): return jsonify({"type": "text", "data": text})

        if tool_type in ["quiz", "flashcard"]:
            json_data = clean_ai_response(text)
            if json_data: return jsonify({"type": "json", "data": json_data})
            else: return jsonify({"type": "text", "data": "Error: AI response was not valid JSON."})
            
        return jsonify({"type": "text", "data": text})
    except Exception as e:
        return jsonify({"type": "text", "data": f"Backend Error: {str(e)}"})


@app.route('/delete-notebook', methods=['POST'])
def delete_notebook():
    notebook_id = request.json.get('id')
    
    nb = get_notebook(notebook_id)
    user_email = nb.get('user_email', 'unknown') if nb else 'unknown'
    title = nb.get('title', 'Unknown Notebook') if nb else 'Unknown Notebook'

    if USING_MONGO:
        notebooks_col.delete_one({"id": notebook_id})
    else:
        global memory_notebooks
        memory_notebooks = [n for n in memory_notebooks if n["id"] != notebook_id]
        
    # ✅ Logging Deletion
    log_activity(user_email, "Deleted Notebook", f"Deleted notebook titled: {title}")
        
    return jsonify({"success": True})

# --- MULTIPLAYER KAHOOT-STYLE ENDPOINTS REMAIN UNCHANGED BELOW ---
@app.route('/host-game', methods=['POST'])
def host_game():
    data = request.json
    code = ''.join(random.choices(string.ascii_uppercase + string.digits, k=6))
    email = data.get('email', 'guest')
    
    active_games[code] = {
        "host_email": email,
        "quiz_data": data.get('quiz_data', []),
        "players": {},
        "status": "waiting"
    }
    
    log_activity(email, "Hosted Live Game", f"Created game room with code: {code}")
    return jsonify({"code": code})

@app.route('/join-game', methods=['POST'])
def join_game():
    data = request.json
    code = data.get('code', '').upper()
    name = data.get('name', 'Anonymous')
    
    if code in active_games:
        if active_games[code]['status'] != 'waiting':
            return jsonify({"error": "Game already started! Too late to join."}), 400
            
        active_games[code]['players'][name] = 0
        log_activity("player_join", "Joined Game", f"Player '{name}' joined room {code}")
        return jsonify({
            "success": True, 
            "quiz_data": active_games[code]['quiz_data']
        })
    return jsonify({"error": "Invalid Game Code"}), 404

@app.route('/start-game', methods=['POST'])
def start_game():
    code = request.json.get('code', '').upper()
    if code in active_games:
        active_games[code]['status'] = 'playing'
        return jsonify({"success": True})
    return jsonify({"error": "Game not found"}), 404

@app.route('/get-game-status', methods=['POST'])
def get_game_status():
    code = request.json.get('code', '').upper()
    if code in active_games:
        players = list(active_games[code]['players'].keys())
        return jsonify({
            "status": active_games[code]['status'],
            "players": players
        })
    return jsonify({"error": "Game not found"}), 404

@app.route('/update-score', methods=['POST'])
def update_score():
    data = request.json
    code = data.get('code', '').upper()
    name = data.get('name')
    score = data.get('score', 0)
    
    if code in active_games and name in active_games[code]['players']:
        active_games[code]['players'][name] = score
        return jsonify({"success": True})
    return jsonify({"error": "Game or player not found"}), 404

@app.route('/get-leaderboard', methods=['POST'])
def get_leaderboard():
    code = request.json.get('code', '').upper()
    if code in active_games:
        players = active_games[code]['players']
        sorted_players = sorted(players.items(), key=lambda x: x[1], reverse=True)
        leaderboard = [{"name": k, "score": v} for k, v in sorted_players]
        return jsonify({"leaderboard": leaderboard, "status": active_games[code]['status']})
    return jsonify({"error": "Game not found"}), 404


if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)